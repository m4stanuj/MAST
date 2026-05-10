#!/usr/bin/env python3
"""
file_mcp.py — OpenWork File Operations MCP Server v2
======================================================
DROP-IN INSTALL (no manual config editing needed):

  python file_mcp.py --install
  → copies self to ~/.config/opencode/
  → patches opencode.json automatically
  → done. Reload workspace.

Supported:
  Text    : .txt .md .log .csv .json .jsonl .yaml .toml .xml .html
            .py .js .ts .jsx .tsx .bat .ps1 .sh .css .sql .ini .env etc.
  Office  : .pdf .docx .xlsx .pptx (text extraction)
  Image   : .png .jpg .jpeg .webp .gif .bmp (info + optional base64)
  Archive : .zip .tar .tar.gz .7z (list / extract / create)
  Binary  : .exe .dll → hex dump + info

Tools (15):
  file_read    file_write    file_info     file_list     file_search
  file_copy    file_move     file_delete
  zip_list     zip_extract   zip_create
  image_info   csv_head      json_peek     office_extract
"""

import sys, os, json, time, shutil, mimetypes, hashlib
from pathlib import Path

_THIS_FILE  = Path(__file__).resolve()
_CONFIG_DIR = Path(os.environ.get("MAST_CONFIG") or os.environ.get("OPENWORK_CONFIG") or  os.path.expanduser("~/.config/opencode"))

# Hardened base
sys.path.insert(0, str(_CONFIG_DIR))
sys.path.insert(0, str(Path(__file__).parent))
from _mcp_base import mcp_loop

def _log(m): print(f"[file_mcp] {m}", file=sys.stderr, flush=True)
def _send(o): print(json.dumps(o, ensure_ascii=False), flush=True)
def _err(rid, msg): _send({"jsonrpc":"2.0","id":rid,"error":{"code":-32000,"message":str(msg)}})

# ─────────────────────────────────────────────────────────────────────
# ANSI HELPERS
# ─────────────────────────────────────────────────────────────────────
_NO_COLOR = not sys.stdout.isatty() or os.environ.get("NO_COLOR")

def _c(code, text):
    return text if _NO_COLOR else f"\033[{code}m{text}\033[0m"

def _bold(t):    return _c("1", t)
def _dim(t):     return _c("2", t)
def _green(t):   return _c("92", t)
def _yellow(t):  return _c("93", t)
def _cyan(t):    return _c("96", t)
def _red(t):     return _c("91", t)
def _magenta(t): return _c("95", t)
def _white(t):   return _c("97", t)
def _blue(t):    return _c("94", t)

# ── Animation engine ─────────────────────────────────────────────────
import math

# Spinner frames — arc sweep feels natural
_SPIN = ["◜","◠","◝","◞","◡","◟"]

def _ease_out_cubic(t: float) -> float:
    """t in [0,1] → eased value. Starts fast, decelerates smoothly."""
    return 1 - (1 - t) ** 3

def _ease_in_out(t: float) -> float:
    """Smooth S-curve — slow start, fast middle, slow end."""
    return 3*t*t - 2*t*t*t

def _gradient_bar(pct: float, width: int = 30) -> str:
    """
    Fills bar with a colour gradient:
      0–40 %  → cyan  ▓
      40–75 % → green █
      75–100% → bright-green █  (leading edge glows)
    Tail chars dim as fill advances for a 'lit tip' effect.
    """
    if _NO_COLOR:
        done = int(width * pct / 100)
        return "[" + "█" * done + "░" * (width - done) + "]"
    cells = []
    filled = width * pct / 100
    for idx in range(width):
        if idx >= filled:
            cells.append(_dim("░"))
        elif idx >= filled - 1:           # leading edge — bright tip
            cells.append(_c("97", "█"))
        elif idx / width < 0.40:
            cells.append(_c("96", "▓"))   # cyan early fill
        elif idx / width < 0.75:
            cells.append(_c("92", "█"))   # green mid fill
        else:
            cells.append(_c("92;1", "█")) # bright-green tail
    return "[" + "".join(cells) + "]"

def _comet_bar(tick: int, width: int = 30) -> str:
    """Indeterminate scanner — a glowing comet bounces left↔right."""
    period = width * 2
    pos    = tick % period
    if pos >= width:
        pos = period - pos
    cells = [_dim("░")] * width
    tail  = [_c("36","▒"), _c("96","▓"), _c("97;1","█"), _c("96","▓"), _c("36","▒")]
    for k, ch in enumerate(tail):
        j = pos - 2 + k
        if 0 <= j < width:
            cells[j] = ch
    return "[" + "".join(cells) + "]"

def _typewrite(line: str, delay: float = 0.018) -> None:
    """Print each char with a tiny delay — natural typing feel."""
    for ch in line:
        print(ch, end="", flush=True)
        time.sleep(delay)
    print()

def _animate_step(label: str, duration: float = 0.5) -> None:
    """
    Eased spinner: starts fast, slows down as it nears completion,
    then snaps to a green ✔.
    """
    steps = 28
    for i in range(steps):
        t     = i / steps
        delay = 0.012 + 0.028 * _ease_in_out(t)   # speeds up then slows
        spin  = _c("96", _SPIN[i % len(_SPIN)])
        trail = _dim("·" * (i % 3 + 1))
        print(f"\r  {spin}  {label}{trail}   ", end="", flush=True)
        time.sleep(delay)
    print(f"\r  {_green('✔')}  {label}   ")

def _animate_bar(label: str, duration: float = 0.7, width: int = 30) -> None:
    """
    Eased progress bar: ease-out-cubic so it rushes in then settles,
    with a glowing gradient fill and a bouncy overshoot feel at 100%.
    """
    fps    = 50
    steps  = int(duration * fps)
    for i in range(steps + 1):
        t      = i / steps
        pct    = _ease_out_cubic(t) * 100
        bar    = _gradient_bar(pct, width)
        pct_lbl = _yellow(f"{int(pct):3d}%")
        print(f"\r  {bar} {pct_lbl}  {_dim(label)}", end="", flush=True)
        time.sleep(1 / fps)
    # snap to crisp 100%
    print(f"\r  {_gradient_bar(100, width)} {_green('100%')}  {_bold(label)}  ")

def _animate_scan(label: str, duration: float = 0.65, width: int = 30) -> None:
    """
    Comet scanner for indeterminate steps (reading / checking).
    Fades out by converting to a full bar at the end.
    """
    fps   = 40
    steps = int(duration * fps)
    for i in range(steps):
        bar = _comet_bar(i, width)
        print(f"\r  {bar}  {_dim(label)}", end="", flush=True)
        time.sleep(1 / fps)
    # morph into solid bar
    _animate_bar(label, duration=0.25, width=width)

def _section(title: str, w: int = 52) -> None:
    """Thin rule with a centered, coloured label — adds breathing room."""
    pad = max(0, w - len(title) - 4)
    l, r = pad // 2, pad - pad // 2
    print(f"\n  {_dim('─' * l)} {_cyan(title)} {_dim('─' * r)}")

def _badge(text: str, color_fn=None) -> str:
    fn = color_fn or _cyan
    return fn(f" {text} ")

# ─────────────────────────────────────────────────────────────────────
# AUTO-INSTALL  →  python file_mcp.py --install
# ─────────────────────────────────────────────────────────────────────
def _auto_install():
    dest     = _CONFIG_DIR / "file_mcp.py"
    cfg_path = _CONFIG_DIR / "opencode.json"
    username = os.environ.get("USERNAME") or os.environ.get("USER") or "user"

    # ── Header banner ──────────────────────────────────────────────
    w = 56
    print()
    print("  " + _cyan("╔" + "═" * w + "╗"))
    print("  " + _cyan("║") + " " * w + _cyan("║"))
    _typewrite("  " + _cyan("║") + "    📁  " + _bold(_white("OpenWork  File MCP")) + _dim("  ·  Installer v2.0") + " " * 14 + _cyan("║"), delay=0.012)
    print("  " + _cyan("║") + " " * w + _cyan("║"))
    print("  " + _cyan("╚" + "═" * w + "╝"))
    print()

    # ── Step 1 ─────────────────────────────────────────────────────
    _section("SETUP")
    time.sleep(0.08)
    _CONFIG_DIR.mkdir(parents=True, exist_ok=True)
    _animate_step(f"Config dir  {_dim(str(_CONFIG_DIR))}")

    # ── Step 2 ─────────────────────────────────────────────────────
    _section("COPY")
    _animate_bar("Copying file_mcp.py", duration=0.8)
    if _THIS_FILE != dest:
        shutil.copy2(_THIS_FILE, dest)
        print(f"  {_green('✔')}  {_bold('Saved')}  →  {_cyan(str(dest))}")
    else:
        print(f"  {_green('✔')}  Already in place: {_dim(str(dest))}")

    # ── Step 3 ─────────────────────────────────────────────────────
    _section("PATCH  opencode.json")
    time.sleep(0.12)

    new_entry = {
        "type":    "local",
        "command": ["python", f"C:/Users/{username}/.config/opencode/file_mcp.py"],
        "enabled": True
    }

    if cfg_path.exists():
        try:
            _animate_scan(f"Reading {_dim('opencode.json')}", duration=0.5)
            cfg = json.loads(cfg_path.read_text(encoding="utf-8"))
        except json.JSONDecodeError as e:
            bak = cfg_path.with_name("opencode.json.bak")
            shutil.copy2(cfg_path, bak)
            print(f"\n  {_yellow('⚠')}  JSON error — backup → {_dim(str(bak))}")
            print(f"  {_red('→')}  {str(e)}")
            cfg = {}
    else:
        cfg = {}
        print(f"  {_yellow('ℹ')}  opencode.json not found — creating fresh")

    cfg.setdefault("mcp", {})
    force = "--force" in sys.argv

    if "file" in cfg["mcp"] and not force:
        _animate_step(f"{_badge('file', _dim)}  already registered — {_yellow('skipped')}", duration=0.25)
        print(f"  {_cyan('ℹ')}  Kept as-is  {_dim('(--force to overwrite)')}")
    else:
        label = f"Force-updating {_badge('file', _cyan)}" if force else f"Injecting {_badge('file', _cyan)} entry"
        _animate_step(label, duration=0.35)
        cfg["mcp"]["file"] = new_entry

    _animate_bar("Writing opencode.json", duration=0.45)
    cfg_path.write_text(json.dumps(cfg, indent=2, ensure_ascii=False), encoding="utf-8")
    print(f"  {_green('✔')}  Saved  →  {_cyan(str(cfg_path))}")

    # ── Deps ───────────────────────────────────────────────────────
    _section("OPTIONAL DEPS")
    deps = [
        ("pypdf",        "📕 PDF extraction"),
        ("python-docx",  "📘 DOCX support"),
        ("openpyxl",     "📗 XLSX support"),
        ("python-pptx",  "📙 PPTX support"),
        ("Pillow",       "🖼️  Image info"),
        ("py7zr",        "🗜️  7z archives"),
    ]
    for pkg, desc in deps:
        print(f"  {_dim('○')}  {_yellow(f'pip install {pkg}'):<32}  {_dim(desc)}")

    # ── Footer ─────────────────────────────────────────────────────
    print()
    print("  " + _cyan("┌" + "─" * w + "┐"))
    print("  " + _cyan("│") + f"  {_bold('🔁  Reload OpenWork workspace to activate.')}" + " " * (w - 44) + _cyan("│"))
    print("  " + _cyan("└" + "─" * w + "┘"))
    print()
    _typewrite(f"  {_green(_bold('✅  Done!  file_mcp is live.'))}  15 tools ready.", delay=0.018)
    print()


# ─────────────────────────────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────────────────────────────
def _resolve(path: str) -> Path:
    p = Path(path).expanduser()
    return p if p.is_absolute() else Path.cwd() / p

def _human_size(b):
    for u in ["B","KB","MB","GB"]:
        if b < 1024: return f"{b:.1f} {u}"
        b /= 1024
    return f"{b:.1f} TB"

def _mime(path: Path):
    mt, _ = mimetypes.guess_type(str(path))
    return mt or "application/octet-stream"

_TEXT_EXTS = {
    ".txt",".md",".log",".csv",".json",".jsonl",".yaml",".yml",
    ".toml",".xml",".html",".htm",".py",".js",".ts",".jsx",".tsx",
    ".bat",".ps1",".sh",".css",".sql",".ini",".cfg",".conf",".env",
    ".rs",".go",".java",".c",".cpp",".h",".cs",".rb",".php",".swift",
    ".kt",".r",".m",".tex",".rst",".nfo",".srt",".vtt",".vue",".svelte"
}
def _is_text(p: Path): return p.suffix.lower() in _TEXT_EXTS

# ─────────────────────────────────────────────────────────────────────
# TOOL HANDLERS
# ─────────────────────────────────────────────────────────────────────
def _t_file_read(a):
    path = _resolve(a.get("path",""))
    if not path.exists(): return f"ERROR: Not found: {path}"
    if not path.is_file(): return f"ERROR: Not a file: {path}"
    ext = path.suffix.lower(); max_c = int(a.get("max_chars",8000))
    if ext in {".pdf",".docx",".xlsx",".pptx"}:
        return _t_office_extract({"path":str(path),"max_chars":max_c})
    if ext in {".png",".jpg",".jpeg",".webp",".gif",".bmp"}:
        return _t_image_info({"path":str(path)})
    if ext==".zip" or any(str(path).endswith(s) for s in (".tar",".tar.gz",".tgz",".tar.bz2",".7z")):
        return _t_zip_list({"path":str(path)})
    if _is_text(path):
        try:
            text=path.read_text(encoding=a.get("encoding","utf-8"),errors="replace")
            lines=text.splitlines(); total=len(lines)
            s=int(a.get("start_line",1)); e=int(a.get("end_line",0)) or total
            chunk="\n".join(lines[s-1:e])
            if len(chunk)>max_c: chunk=chunk[:max_c]+f"\n\n...[truncated — {total} lines total]"
            return f"📄 {path.name} ({total} lines, {_human_size(path.stat().st_size)})\n\n{chunk}"
        except Exception as ex: return f"ERROR: {ex}"
    try:
        raw=path.read_bytes()[:256]
        return f"Binary: {path.name} ({_human_size(path.stat().st_size)})\nHex:\n{' '.join(f'{b:02x}' for b in raw)}"
    except Exception as ex: return f"ERROR: {ex}"

def _t_file_write(a):
    path=_resolve(a.get("path","")); content=a.get("content","")
    mode=a.get("mode","write"); enc=a.get("encoding","utf-8")
    try:
        path.parent.mkdir(parents=True,exist_ok=True)
        if mode=="append":
            with open(path,"a",encoding=enc) as f: f.write(content)
            return f"✅ Appended {len(content)} chars → {path}"
        path.write_text(content,encoding=enc)
        return f"✅ Written {len(content)} chars → {path}"
    except Exception as ex: return f"ERROR: {ex}"

def _t_file_info(a):
    path=_resolve(a.get("path",""))
    if not path.exists(): return f"ERROR: Not found: {path}"
    st=path.stat()
    lines=[f"{'📁' if path.is_dir() else '📄'} {path.name}",
           f"  path:     {path}",
           f"  type:     {'directory' if path.is_dir() else 'file'}",
           f"  size:     {_human_size(st.st_size)}",
           f"  mime:     {_mime(path) if path.is_file() else 'inode/directory'}",
           f"  modified: {time.strftime('%Y-%m-%d %H:%M:%S',time.localtime(st.st_mtime))}",
           f"  created:  {time.strftime('%Y-%m-%d %H:%M:%S',time.localtime(st.st_ctime))}",
           f"  ext:      {path.suffix.lower()}"]
    if path.is_file() and _is_text(path):
        try:
            txt=path.read_text(encoding="utf-8",errors="replace")
            lines+=[f"  lines:    {len(txt.splitlines())}",f"  chars:    {len(txt)}"]
        except: pass
    if a.get("include_hash") and path.is_file():
        try: lines.append(f"  md5:      {hashlib.md5(path.read_bytes()).hexdigest()}")
        except: pass
    return "\n".join(lines)

def _t_file_list(a):
    path=_resolve(a.get("path","."))
    if not path.exists(): return f"ERROR: Not found: {path}"
    if not path.is_dir(): return f"ERROR: Not a directory: {path}"
    pat=a.get("pattern","*"); show=a.get("show","all"); lim=int(a.get("limit",100))
    try:
        items=list(path.rglob(pat) if a.get("recursive") else path.glob(pat))
        if show=="files": items=[i for i in items if i.is_file()]
        elif show=="dirs": items=[i for i in items if i.is_dir()]
        items.sort(key=lambda x:(x.is_file(),x.name.lower()))
        total=len(items); items=items[:lim]
        lines=[f"📂 {path}  ({total} items{', first '+str(lim) if total>lim else ''})"]
        for it in items:
            rel=it.relative_to(path)
            sz=f"  ({_human_size(it.stat().st_size)})" if it.is_file() else ""
            lines.append(f"  {'📄' if it.is_file() else '📁'} {rel}{'/' if it.is_dir() else ''}{sz}")
        return "\n".join(lines)
    except Exception as ex: return f"ERROR: {ex}"

def _t_file_search(a):
    path=_resolve(a.get("path",".")); query=a.get("query","")
    if not query: return "ERROR: query required"
    pat=a.get("pattern","*"); case=bool(a.get("case_sensitive",False)); lim=int(a.get("limit",50))
    results=[]
    try:
        files=list(path.rglob(pat)) if path.is_dir() else [path]
        for f in files:
            if not f.is_file() or not _is_text(f): continue
            try:
                for i,line in enumerate(f.read_text(encoding="utf-8",errors="replace").splitlines(),1):
                    if (query if case else query.lower()) in (line if case else line.lower()):
                        results.append(f"  {f}:{i}  {line.strip()[:120]}")
                        if len(results)>=lim: break
            except: continue
            if len(results)>=lim: break
        return (f"🔍 {len(results)} matches for '{query}':\n"+"\n".join(results)) if results else f"No matches for '{query}'"
    except Exception as ex: return f"ERROR: {ex}"

def _t_file_copy(a):
    src=_resolve(a.get("src","")); dest=_resolve(a.get("dest",""))
    if not src.exists(): return f"ERROR: Not found: {src}"
    try:
        dest.parent.mkdir(parents=True,exist_ok=True)
        if src.is_dir():
            shutil.copytree(src, dest, dirs_exist_ok=True)
        else:
            shutil.copy2(src, dest)
        return f"✅ Copied {src} → {dest}"
    except Exception as ex: return f"ERROR: {ex}"

def _t_file_move(a):
    src=_resolve(a.get("src","")); dest=_resolve(a.get("dest",""))
    if not src.exists(): return f"ERROR: Not found: {src}"
    try:
        dest.parent.mkdir(parents=True,exist_ok=True)
        shutil.move(str(src),str(dest)); return f"✅ Moved {src} → {dest}"
    except Exception as ex: return f"ERROR: {ex}"

def _t_file_delete(a):
    path=_resolve(a.get("path",""))
    if not path.exists(): return f"ERROR: Not found: {path}"
    if not a.get("confirm"): return f"⚠️ Add confirm=true to delete: {path}"
    try:
        if path.is_dir():
            shutil.rmtree(path)
        else:
            path.unlink(missing_ok=True)
        return f"🗑️ Deleted: {path}"
    except Exception as ex: return f"ERROR: {ex}"

def _t_zip_list(a):
    path=_resolve(a.get("path",""))
    if not path.exists(): return f"ERROR: Not found: {path}"
    ext=path.suffix.lower()
    try:
        if ext==".zip":
            import zipfile
            with zipfile.ZipFile(path) as z:
                ii=z.infolist(); sz=sum(i.file_size for i in ii)
                lines=[f"🗜️ {path.name}  ({len(ii)} files, {_human_size(sz)} uncompressed)"]
                for i in ii[:80]: lines.append(f"  {i.filename}  ({_human_size(i.file_size)})")
                if len(ii)>80: lines.append(f"  ...+{len(ii)-80} more")
                return "\n".join(lines)
        elif any(str(path).endswith(s) for s in (".tar",".tar.gz",".tgz",".tar.bz2",".tar.xz")):
            import tarfile
            with tarfile.open(path) as t:
                mm=t.getmembers(); lines=[f"🗜️ {path.name}  ({len(mm)} entries)"]
                for m in mm[:80]: lines.append(f"  {m.name}  ({_human_size(m.size)})")
                if len(mm)>80: lines.append(f"  ...+{len(mm)-80} more")
                return "\n".join(lines)
        elif ext==".7z":
            import py7zr
            with py7zr.SevenZipFile(path,mode='r') as z:
                nn=z.getnames(); lines=[f"🗜️ {path.name}  ({len(nn)} files)"]
                for n in nn[:80]: lines.append(f"  {n}")
                if len(nn)>80: lines.append(f"  ...+{len(nn)-80} more")
                return "\n".join(lines)
        return f"ERROR: Unsupported archive: {ext}"
    except ImportError as e: return f"ERROR: Missing dep — {e}. Run: pip install py7zr"
    except Exception as ex: return f"ERROR: {ex}"

def _t_zip_extract(a):
    path=_resolve(a.get("path",""))
    # Strip ALL suffixes for multi-ext archives like .tar.gz → "archive" not "archive.tar"
    def _clean_stem(p):
        s = p.name
        for ext in (".tar.gz",".tar.bz2",".tar.xz",".tgz"):
            if s.endswith(ext): return s[:-len(ext)]
        return p.stem
    dest=_resolve(a.get("dest", str(path.parent / _clean_stem(path)) if path.suffix else str(path.parent/"extracted")))
    member=a.get("member","")
    if not path.exists(): return f"ERROR: Not found: {path}"
    ext=path.suffix.lower()
    try:
        dest.mkdir(parents=True,exist_ok=True)
        if ext==".zip":
            import zipfile
            with zipfile.ZipFile(path) as z:
                if member: z.extract(member,dest); return f"✅ Extracted '{member}' → {dest}"
                z.extractall(dest); return f"✅ Extracted {len(z.namelist())} files → {dest}"
        elif any(str(path).endswith(s) for s in (".tar",".tar.gz",".tgz",".tar.bz2",".tar.xz")):
            import tarfile
            with tarfile.open(path) as t:
                if member: t.extract(member,dest); return f"✅ Extracted '{member}' → {dest}"
                t.extractall(dest); return f"✅ Extracted {len(t.getnames())} files → {dest}"
        elif ext==".7z":
            import py7zr
            with py7zr.SevenZipFile(path,mode='r') as z:
                z.extract(dest,targets=[member] if member else None)
                return f"✅ Extracted → {dest}"
        return f"ERROR: Unsupported archive: {ext}"
    except ImportError as e: return f"ERROR: {e}. Run: pip install py7zr"
    except Exception as ex: return f"ERROR: {ex}"

def _t_zip_create(a):
    out=_resolve(a.get("output","")); srcs=[_resolve(s) for s in a.get("sources",[])]
    if not out: return "ERROR: output required"
    if not srcs: return "ERROR: sources required"
    try:
        import zipfile; count=0
        with zipfile.ZipFile(out,"w",zipfile.ZIP_DEFLATED) as z:
            for src in srcs:
                if src.is_dir():
                    for f in src.rglob("*"):
                        if f.is_file(): z.write(f,f.relative_to(src.parent)); count+=1
                elif src.is_file(): z.write(src,src.name); count+=1
        return f"✅ Created {out}  ({count} files, {_human_size(out.stat().st_size)})"
    except Exception as ex: return f"ERROR: {ex}"

def _t_image_info(a):
    path=_resolve(a.get("path",""))
    if not path.exists(): return f"ERROR: Not found: {path}"
    try:
        from PIL import Image
        with Image.open(path) as img:
            lines=[f"🖼️ {path.name}",f"  format: {img.format}",
                   f"  size:   {img.width} × {img.height} px",
                   f"  mode:   {img.mode}",f"  bytes:  {_human_size(path.stat().st_size)}"]
            if a.get("include_base64"):
                import base64; b64=base64.b64encode(path.read_bytes()).decode()
                lines.append(f"\nbase64[:200]: {b64[:200]}...")
            return "\n".join(lines)
    except ImportError: return f"🖼️ {path.name} ({_human_size(path.stat().st_size)}) — pip install Pillow"
    except Exception as ex: return f"ERROR: {ex}"

def _t_csv_head(a):
    path=_resolve(a.get("path","")); rows=int(a.get("rows",10))
    if not path.exists(): return f"ERROR: Not found: {path}"
    try:
        import csv
        all_rows = []
        with open(path,encoding="utf-8",errors="replace") as f:
            reader=csv.reader(f)
            for i, row in enumerate(reader):
                all_rows.append(row)
                if i >= rows: break  # header + rows lines
        if not all_rows: return "Empty CSV"
        # Count total lines separately (fast)
        with open(path, encoding="utf-8", errors="replace") as f:
            total = max(0, sum(1 for _ in f) - 1)  # subtract header
        header = all_rows[0]
        data   = all_rows[1:]
        sep=" | "
        out=[sep.join(str(c)[:30] for c in header)]
        out.append("-"*min(120,len(out[0])))
        for row in data: out.append(sep.join(str(c)[:30] for c in row))
        return f"CSV: {path.name} ({total} rows, {len(header)} cols)\n\n"+"\n".join(out)
    except Exception as ex: return f"ERROR: {ex}"

def _t_json_peek(a):
    path=_resolve(a.get("path",""))
    if not path.exists(): return f"ERROR: Not found: {path}"
    try:
        text=path.read_text(encoding="utf-8",errors="replace")
        if path.suffix.lower()==".jsonl":
            lines=[l.strip() for l in text.splitlines() if l.strip()]
            if lines:
                first=json.loads(lines[0])
                return (f"JSONL: {path.name} ({len(lines)} lines)\n"
                        f"Keys: {list(first.keys()) if isinstance(first,dict) else type(first).__name__}\n"
                        f"Sample: {json.dumps(first,ensure_ascii=False)[:500]}")
        data=json.loads(text)
        if isinstance(data,dict): return f"JSON object: {len(data)} keys\nKeys: {list(data.keys())[:30]}\nPreview: {json.dumps(data,ensure_ascii=False)[:600]}"
        elif isinstance(data,list):
            first=data[0] if data else {}
            return (f"JSON array: {len(data)} items\n"
                    f"First: {list(first.keys()) if isinstance(first,dict) else type(first).__name__}\n"
                    f"Sample: {json.dumps(first,ensure_ascii=False)[:500]}")
        return f"JSON scalar: {type(data).__name__} = {str(data)[:200]}"
    except Exception as ex: return f"ERROR: {ex}"

def _t_office_extract(a):
    path=_resolve(a.get("path","")); max_c=int(a.get("max_chars",8000))
    if not path.exists(): return f"ERROR: Not found: {path}"
    ext=path.suffix.lower()
    if ext==".pdf":
        try:
            from pypdf import PdfReader; r=PdfReader(str(path)); text=""
            for pg in r.pages:
                text+=pg.extract_text() or ""
                if len(text)>max_c: break
            return f"📕 PDF: {path.name} ({len(r.pages)} pages)\n\n{text[:max_c]}"
        except ImportError: return "ERROR: pip install pypdf"
        except Exception as ex: return f"ERROR: {ex}"
    elif ext==".docx":
        try:
            from docx import Document; doc=Document(str(path))
            text="\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return f"📘 DOCX: {path.name} ({len(doc.paragraphs)} paragraphs)\n\n{text[:max_c]}"
        except ImportError: return "ERROR: pip install python-docx"
        except Exception as ex: return f"ERROR: {ex}"
    elif ext==".xlsx":
        try:
            import openpyxl; wb=openpyxl.load_workbook(str(path),read_only=True,data_only=True)
            lines=[f"📗 XLSX: {path.name} ({len(wb.sheetnames)} sheets: {', '.join(wb.sheetnames)})"]
            for sn in wb.sheetnames[:3]:
                ws=wb[sn]; lines.append(f"\n[{sn}]"); count=0
                for row in ws.iter_rows(values_only=True):
                    if all(c is None for c in row): continue
                    lines.append("  "+" | ".join(str(c)[:25] if c is not None else "" for c in row)); count+=1
                    if count>=20: lines.append("  ..."); break
            return "\n".join(lines)[:max_c]
        except ImportError: return "ERROR: pip install openpyxl"
        except Exception as ex: return f"ERROR: {ex}"
    elif ext==".pptx":
        try:
            from pptx import Presentation; prs=Presentation(str(path))
            lines=[f"📙 PPTX: {path.name} ({len(prs.slides)} slides)"]
            for i,sl in enumerate(prs.slides,1):
                texts=[sh.text for sh in sl.shapes if hasattr(sh,"text") and sh.text.strip()]
                lines.append(f"\n[Slide {i}]  "+"  |  ".join(t[:80] for t in texts[:4]))
                if len("\n".join(lines))>max_c: lines.append("..."); break
            return "\n".join(lines)
        except ImportError: return "ERROR: pip install python-pptx"
        except Exception as ex: return f"ERROR: {ex}"
    return f"ERROR: Unsupported: {ext}"

# ─────────────────────────────────────────────────────────────────────
# REGISTRY
# ─────────────────────────────────────────────────────────────────────
TOOLS = {
    "file_read":      (_t_file_read,      "Read any file — auto-detects text/office/image/archive/binary."),
    "file_write":     (_t_file_write,     "Write or append to a text file. Creates directories automatically."),
    "file_info":      (_t_file_info,      "File metadata: size, type, mime, dates, line count, optional md5."),
    "file_list":      (_t_file_list,      "List directory. Supports glob pattern and recursive mode."),
    "file_search":    (_t_file_search,    "Search text inside files. Recursive glob supported."),
    "file_copy":      (_t_file_copy,      "Copy file or folder to destination."),
    "file_move":      (_t_file_move,      "Move or rename file or folder."),
    "file_delete":    (_t_file_delete,    "Delete file/folder. Requires confirm=true for safety."),
    "zip_list":       (_t_zip_list,       "List .zip / .tar / .tar.gz / .7z contents."),
    "zip_extract":    (_t_zip_extract,    "Extract archive. Optional single-member extraction."),
    "zip_create":     (_t_zip_create,     "Create .zip from files/folders."),
    "image_info":     (_t_image_info,     "Image dimensions, format, mode. Optional base64 output."),
    "csv_head":       (_t_csv_head,       "Preview first N rows of CSV as table."),
    "json_peek":      (_t_json_peek,      "Inspect JSON/JSONL — keys, length, sample data."),
    "office_extract": (_t_office_extract, "Extract text from PDF / DOCX / XLSX / PPTX."),
}
SCHEMAS = {
    "file_read":      {"type":"object","properties":{"path":{"type":"string"},"max_chars":{"type":"integer"},"start_line":{"type":"integer"},"end_line":{"type":"integer"},"encoding":{"type":"string"},"include_base64":{"type":"boolean"}},"required":["path"]},
    "file_write":     {"type":"object","properties":{"path":{"type":"string"},"content":{"type":"string"},"mode":{"type":"string","enum":["write","append"]},"encoding":{"type":"string"}},"required":["path","content"]},
    "file_info":      {"type":"object","properties":{"path":{"type":"string"},"include_hash":{"type":"boolean"}},"required":["path"]},
    "file_list":      {"type":"object","properties":{"path":{"type":"string"},"pattern":{"type":"string"},"recursive":{"type":"boolean"},"show":{"type":"string","enum":["all","files","dirs"]},"limit":{"type":"integer"}},"required":["path"]},
    "file_search":    {"type":"object","properties":{"path":{"type":"string"},"query":{"type":"string"},"pattern":{"type":"string"},"case_sensitive":{"type":"boolean"},"limit":{"type":"integer"}},"required":["path","query"]},
    "file_copy":      {"type":"object","properties":{"src":{"type":"string"},"dest":{"type":"string"}},"required":["src","dest"]},
    "file_move":      {"type":"object","properties":{"src":{"type":"string"},"dest":{"type":"string"}},"required":["src","dest"]},
    "file_delete":    {"type":"object","properties":{"path":{"type":"string"},"confirm":{"type":"boolean"}},"required":["path"]},
    "zip_list":       {"type":"object","properties":{"path":{"type":"string"}},"required":["path"]},
    "zip_extract":    {"type":"object","properties":{"path":{"type":"string"},"dest":{"type":"string"},"member":{"type":"string"}},"required":["path"]},
    "zip_create":     {"type":"object","properties":{"output":{"type":"string"},"sources":{"type":"array","items":{"type":"string"}}},"required":["output","sources"]},
    "image_info":     {"type":"object","properties":{"path":{"type":"string"},"include_base64":{"type":"boolean"}},"required":["path"]},
    "csv_head":       {"type":"object","properties":{"path":{"type":"string"},"rows":{"type":"integer"}},"required":["path"]},
    "json_peek":      {"type":"object","properties":{"path":{"type":"string"}},"required":["path"]},
    "office_extract": {"type":"object","properties":{"path":{"type":"string"},"max_chars":{"type":"integer"}},"required":["path"]},
}

# ─────────────────────────────────────────────────────────────────────
# MCP LOOP
# ─────────────────────────────────────────────────────────────────────
def _handle(req):
    m,rid=req.get("method",""),req.get("id")
    if m=="initialize":
        _send({"jsonrpc":"2.0","id":rid,"result":{"protocolVersion":"2024-11-05","capabilities":{"tools":{}},"serverInfo":{"name":"file-mcp","version":"2.0.0"}}})
    elif m=="tools/list":
        _send({"jsonrpc":"2.0","id":rid,"result":{"tools":[{"name":n,"description":fd[1],"inputSchema":SCHEMAS[n]} for n,fd in TOOLS.items()]}})
    elif m=="tools/call":
        p=req.get("params",{}); tn,args=p.get("name",""),p.get("arguments",{})
        if tn in TOOLS:
            try: _send({"jsonrpc":"2.0","id":rid,"result":{"content":[{"type":"text","text":str(TOOLS[tn][0](args))}]}})
            except Exception as ex: _err(rid,str(ex))
        else: _err(rid,f"Tool not found: {tn}")
    elif m=="notifications/initialized": pass
    elif rid is not None: _err(rid,f"Unknown: {m}")

def main():
    _log("📁 file-mcp v2.0 started — 15 tools ready")
    mcp_loop("file", _handle)

if __name__=="__main__":
    if "--install" in sys.argv:
        _auto_install()
    else:
        main()
